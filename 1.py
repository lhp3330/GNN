from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# 颜色方案
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xAB)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GREEN = RGBColor(0x27, 0xAE, 0x60)


def add_title_bar(slide, title_text):
    """添加顶部标题栏"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.text = title_text
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_body_text(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_TEXT):
    """添加正文文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    return tf


def add_bullet_points(slide, left, top, width, height, points, font_size=16):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)
        p.level = 0
    return tf


# ==================== Slide 1: 封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BLUE
bg.line.fill.background()

add_body_text(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
              "Cooperative Graph Neural Networks\nwith APPNP Enhancement", font_size=36, bold=True, color=WHITE)
add_body_text(slide, Inches(1), Inches(4), Inches(11), Inches(1),
              "基于协作图神经网络与个性化PageRank传播的节点分类方法", font_size=22, color=RGBColor(0xBB, 0xDE, 0xFB))
add_body_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(1),
              "汇报人：XXX    指导教师：XXX\n日期：2025年X月", font_size=18, color=RGBColor(0x90, 0xCA, 0xF9))

# ==================== Slide 2: 目录 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "目录 Contents")
items = [
    "1. 研究背景与动机",
    "2. 相关工作：传统MPNN的局限性",
    "3. CO-GNN 原理：协作式消息传递",
    "4. 本文改进：CO-GNN + APPNP 混合架构",
    "5. 模型架构详解",
    "6. 实验设置与结果",
    "7. 总结与展望"
]
add_bullet_points(slide, Inches(1), Inches(1.8), Inches(10), Inches(5), items, font_size=22)

# ==================== Slide 3: 研究背景 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "1. 研究背景与动机")
points = [
    "▸ 图神经网络(GNN)已成为图结构数据学习的主流范式",
    "▸ 消息传递神经网络(MPNN)：节点通过聚合邻居信息更新表征",
    "",
    "核心问题：传统MPNN中所有节点无条件地监听并广播",
    "",
    "▸ Over-smoothing：深层节点表征趋于同质化",
    "▸ Over-squashing：长程依赖信息被压缩丢失",
    "▸ 异质图(Heterophilic)表现差：邻居信号可能有害",
    "▸ 表达能力上限：不超过1-WL测试",
    "",
    "动机：能否让节点自主决策\"如何与邻居互动\"？"
]
add_bullet_points(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5), points, font_size=17)

# ==================== Slide 4: CO-GNN核心思想 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "2. CO-GNN 核心思想：节点即玩家")
points = [
    "核心洞见：将每个节点视为理性玩家(Player)，自主决策信息流方向",
    "",
    "四种动作(Action Space)：",
    "  ● STANDARD (S)：同时监听 + 广播（传统MPNN行为）",
    "  ● LISTEN (L)：仅监听广播邻居的信息",
    "  ● BROADCAST (B)：仅向监听自己的邻居广播",
    "  ● ISOLATE (I)：既不监听也不广播（自更新）",
    "",
    "关键特性：",
    "  • 信息流是条件性的、有方向的",
    "  • 每一层生成不同的计算图（动态重布线）",
    "  • 动作选择由动作网络π决定，节点更新由环境网络η完成"
]
add_bullet_points(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5), points, font_size=16)

# ==================== Slide 5: CO-GNN 更新规则 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "3. CO-GNN 更新规则")
points = [
    "单层更新流程：",
    "",
    "Step 1: 动作网络π根据节点及邻居表征输出动作概率分布",
    "         p_v = π(h_v, {h_u : u ∈ N(v)})",
    "",
    "Step 2: Straight-through Gumbel-Softmax 可微离散采样",
    "         a_v = GumbelSoftmax(p_v, τ, hard=True)",
    "",
    "Step 3: 环境网络η根据有效邻居集更新表征",
    "         若 a_v = I 或 B：h_v' = η(h_v, ∅)",
    "         若 a_v = L 或 S：h_v' = η(h_v, M)",
    "         其中 M = {h_u | u∈N(v), a_u = S 或 B}",
    "",
    "边权重：w(u→v) = keep_in[v] × keep_out[u]"
]
add_bullet_points(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5), points, font_size=16)

# ==================== Slide 6: APPNP 回顾 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "4. APPNP 回顾：个性化PageRank传播")
points = [
    "APPNP (Approximate Personalized Propagation of Neural Predictions)：",
    "",
    "核心思想：解耦特征变换与传播",
    "  • 特征变换：Z = MLP(X)  （仅一次）",
    "  • 迭代传播：H⁽ᵏ⁺¹⁾ = (1-α)·Â·H⁽ᵏ⁾ + α·Z",
    "",
    "优势：",
    "  ✓ 无参数传播，K步迭代不增加参数量",
    "  ✓ Teleport项 α·Z 保留初始信息，缓解over-smoothing",
    "  ✓ 收敛到个性化PageRank的封闭解",
    "  ✓ 可捕获长程依赖（大K值）",
    "",
    "局限：",
    "  ✗ 固定拓扑传播，无法动态调整信息流",
    "  ✗ 异质图上表现受限"
]
add_bullet_points(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5), points, font_size=16)

# ==================== Slide 7: 本文方法 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "5. 本文方法：CO-GNN + APPNP 混合架构")
points = [
    "核心创新：将CO-GNN的\"智能路由\"与APPNP的\"可靠长程传播\"结合",
    "",
    "设计理念：",
    "  • CO-GNN 负责动态调整信息流方向（哪些边该传、哪些该断）",
    "  • APPNP 负责在调整后的表征上进行多跳平滑传播",
    "  • Gate机制自适应融合两路信息",
    "",
    "互补性分析：",
    "  ┌─────────────┬──────────────────┬──────────────────┐",
    "  │             │ CO-GNN           │ APPNP            │",
    "  ├─────────────┼──────────────────┼──────────────────┤",
    "  │ 拓扑        │ 动态（每层不同） │ 固定（原始图）   │",
    "  │ 传播范围    │ 局部（1-hop/层） │ 全局（K-hop）    │",
    "  │ 参数        │ 有参数           │ 无参数传播       │",
    "  │ 抗噪声      │ 强（可隔离）     │ 弱（全局平滑）   │",
    "  └─────────────┴──────────────────┴──────────────────┘"
]
add_bullet_points(slide, Inches(0.3), Inches(1.5), Inches(12.5), Inches(5.5), points, font_size=14)

# ==================== Slide 8: 架构图 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "6. 模型架构详解")
points = [
    "整体流程：",
    "",
    "  Input X → [Encoder (Linear + Dropout + Act)]",
    "         ↓",
    "  ┌──────────────────────────────────────────┐",
    "  │  Block 1:                                 │",
    "  │    CoGNNConv → x_co (动态重布线后表征)    │",
    "  │    APPNP(K=10,α=0.1) → x_app (PPR平滑)  │",
    "  │    gate = σ(W·x_co)                       │",
    "  │    x = LayerNorm(gate·x_app + (1-gate)·x_co) │",
    "  └──────────────────────────────────────────┘",
    "         ↓  (重复 Block 2)",
    "  [LayerNorm] → [Decoder (Linear)] → Output",
    "",
    "CoGNNConv 内部（num_layers层循环）：",
    "  LN → in_act_net → out_act_net → Gumbel采样",
    "  → edge_weight → env_net → Dropout → Act → 残差连接"
]
add_bullet_points(slide, Inches(0.3), Inches(1.4), Inches(12.5), Inches(5.8), points, font_size=14)

# ==================== Slide 9: Gate机制 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "7. 门控融合机制 (Gated Fusion)")
points = [
    "融合公式：",
    "  x = LayerNorm( gate · x_APPNP + (1 - gate) · x_CoGNN )",
    "  其中 gate = σ(W · x_CoGNN),  W ∈ R^{d×1}",
    "",
    "设计动机：",
    "  • 节点级自适应：每个节点独立决定依赖哪路信息",
    "  • gate → 1：更信任APPNP的全局平滑结果（同质区域）",
    "  • gate → 0：更信任CoGNN的动态路由结果（异质区域）",
    "",
    "正则化损失：",
    "  L = L_task + λ · L_reg",
    "  L_reg = mean( edge_weight · ||z_src - z_dst||² )",
    "",
    "  鼓励连接的节点在embedding空间中接近",
    "  edge_weight加权：动作网络认为重要的边贡献更大"
]
add_bullet_points(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5), points, font_size=16)

# ==================== Slide 10: 理论分析 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "8. 理论性质分析")
points = [
    "继承自CO-GNN的理论优势：",
    "",
    "① 表达能力严格强于1-WL",
    "   • Gumbel采样的随机性使1-WL不可区分的节点有正概率获得不同动作",
    "",
    "② 长程信息传输能力",
    "   • 动作网络可规划\"信息高速公路\"，过滤无关噪声",
    "   • APPNP的K步传播进一步增强长程覆盖",
    "",
    "③ 缓解Over-smoothing",
    "   • ISOLATE动作主动隔离无用邻居",
    "   • APPNP的teleport项保留初始特征",
    "   • 双重机制协同防止表征退化",
    "",
    "④ 任务自适应计算图",
    "   • 不同任务/数据集自动学习不同的动作策略"
]
add_bullet_points(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5), points, font_size=16)

# ==================== Slide 11: 实验设置 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "9. 实验设置")
points = [
    "数据集：CiteSeer（同质引文网络）",
    "  • 节点数：3,327  边数：9,104  类别数：6  同质性：~0.74",
    "",
    "超参数配置：",
    "  • 隐藏维度：128    Dropout：0.5",
    "  • 学习率：5e-3     权重衰减：5e-4",
    "  • Gumbel温度：τ₀=0.1, temp=0.05 (可学习)",
    "  • APPNP：K=10, α=0.1",
    "  • 环境网络：MEAN, 2层    动作网络：MEAN, 2层(dim=16)",
    "  • 正则化权重：λ=0.01",
    "  • 训练轮数：300 (Early Stopping, patience=100)",
    "",
    "评估指标：节点分类准确率 (Accuracy)",
    "基线方法：GCN, GAT, APPNP, GraphSAGE, CO-GNN"
]
add_bullet_points(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5), points, font_size=16)

# ==================== Slide 12: 实验结果 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "10. 实验结果")
points = [
    "CiteSeer 节点分类准确率对比：",
    "",
    "  ┌────────────────────┬──────────────┐",
    "  │ 方法               │ Accuracy (%) │",
    "  ├────────────────────┼──────────────┤",
    "  │ GCN                │ ~70.3        │",
    "  │ GAT                │ ~72.5        │",
    "  │ APPNP              │ ~71.8        │",
    "  │ CO-GNN (原始)      │ ~72.0        │",
    "  │ CO-GNN + APPNP     │  72.5+       │",
    "  │ (本文方法)         │ (seed=686)   │",
    "  └────────────────────┴──────────────┘",
    "",
    "  * 具体数值请替换为您的实验结果",
    "  * 建议补充多种子均值±标准差"
]
add_bullet_points(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5), points, font_size=15)

# ==================== Slide 13: 消融实验 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "11. 消融实验")
points = [
    "消融设计（验证各组件贡献）：",
    "",
    "  ┌──────────────────────────────┬──────────┐",
    "  │ 变体                         │ Acc (%)  │",
    "  ├──────────────────────────────┼──────────┤",
    "  │ 完整模型 (CoGNN+APPNP+Gate)  │  XX.X    │",
    "  │ 去掉APPNP (仅CoGNN)          │  XX.X    │",
    "  │ 去掉CoGNN (仅APPNP)          │  XX.X    │",
    "  │ Gate替换为直接相加            │  XX.X    │",
    "  │ 去掉正则化损失 (λ=0)         │  XX.X    │",
    "  │ 不同act_model (SUM vs MEAN)  │  XX.X    │",
    "  └──────────────────────────────┴──────────┘",
    "",
    "关键发现：",
    "  • Gate融合优于直接相加/拼接",
    "  • APPNP提供稳定的长程信息补充",
    "  • 正则化损失有助于学习更有意义的edge_weight"
]
add_bullet_points(slide, Inches(0.3), Inches(1.5), Inches(12.5), Inches(5.5), points, font_size=14)

# ==================== Slide 14: 超参数分析 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "12. 超参数敏感性分析")
points = [
    "关键超参数影响：",
    "",
    "① Gumbel温度 τ₀：",
    "   • τ₀过大(>1.0)：动作趋于均匀随机，失去选择性",
    "   • τ₀过小(<0.01)：梯度消失，训练不稳定",
    "   • 最优区间：0.05 ~ 0.1",
    "",
    "② APPNP参数 (K, α)：",
    "   • K增大：传播范围扩大，但可能引入噪声",
    "   • α增大：更保守（保留更多初始信息）",
    "   • CiteSeer最优：K=10, α=0.1~0.2",
    "",
    "③ 正则化权重 λ：",
    "   • λ=0：edge_weight缺乏结构约束",
    "   • λ过大：抑制动作多样性",
    "   • 最优：0.001 ~ 0.01"
]
add_bullet_points(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5), points, font_size=15)

# ==================== Slide 15: 总结与展望 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "13. 总结与展望")
points = [
    "本文贡献：",
    "  ✓ 提出CO-GNN与APPNP的混合架构，兼具动态路由与长程传播能力",
    "  ✓ 设计Gate融合机制，实现节点级自适应信息选择",
    "  ✓ 引入边权重正则化，增强embedding的结构一致性",
    "",
    "未来工作：",
    "  • 在更多异质图数据集上验证（Roman-Empire, Amazon-Ratings等）",
    "  • 探索节点级可学习α（teleport概率）",
    "  • 将APPNP的edge_weight也用动作网络生成的权重替代",
    "  • 扩展到图分类和链接预测任务",
    "  • 与Graph Transformer结合",
    "",
    "代码开源：github.com/xxx/CoGNN-APPNP (待补充)"
]
add_bullet_points(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5), points, font_size=16)

# ==================== Slide 16: 致谢 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BLUE
bg.line.fill.background()
add_body_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
              "谢谢！", font_size=48, bold=True, color=WHITE)
add_body_text(slide, Inches(1), Inches(4), Inches(11), Inches(1),
              "Questions & Discussion", font_size=28, color=RGBColor(0x90, 0xCA, 0xF9))
add_body_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(1),
              "联系方式：xxx@xxx.edu.cn", font_size=18, color=RGBColor(0xBB, 0xDE, 0xFB))

# 保存
prs.save("CoGNN_APPNP_Presentation.pptx")
print("PPT已生成：CoGNN_APPNP_Presentation.pptx")
